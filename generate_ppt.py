import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Define color scheme
    bg_color = RGBColor(15, 23, 42)      # Slate 900
    title_color = RGBColor(99, 102, 241)  # Indigo 500
    text_color = RGBColor(248, 250, 252)  # Slate 50
    muted_color = RGBColor(148, 163, 184) # Slate 400

    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = title_color
        return title_box

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6] # blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "InterviewAce AI"
    p.font.name = "Arial"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Interview Trainer Agent using IBM watsonx.ai"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.color.rgb = title_color
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf.add_paragraph()
    p3.text = "\nInternship Submission Project | Powered by IBM Granite LLM"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.color.rgb = muted_color
    p3.alignment = PP_ALIGN.CENTER

    # Helper for bullet points
    def add_bullets_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_title(slide, title)
        
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for idx, item in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item
            p.font.name = "Arial"
            p.font.size = Pt(20)
            p.font.color.rgb = text_color
            p.space_after = Pt(15)
            # Check if sub-bullet
            if item.startswith("  -") or item.startswith("    ") or item.startswith("•"):
                p.level = 1
                p.font.size = Pt(18)
                p.font.color.rgb = muted_color
        return slide

    # Slide 2: Problem Statement
    add_bullets_slide("Problem Statement", [
        "Traditional interview preparation is passive (reading Q&A lists) rather than interactive.",
        "Candidates lack access to instantaneous, quality, expert-level feedback on their explanations.",
        "Hiring a human coach or mentor is highly expensive and not scalable for many job seekers.",
        "Problem Statement No. 22: 'Interview Trainer Agent' - requires an automated, AI-powered agent to mock-interview candidates and grade them objectively."
    ])

    # Slide 3: Existing System
    add_bullets_slide("Existing System & Limitations", [
        "Static Websites: Offer standard list of popular questions without personalized evaluation.",
        "Peer-to-Peer Mocks: Hard to schedule, highly subjective, and dependent on the peer's own level of expertise.",
        "Lack of Scoring: Manual revision does not provide quantitative metrics or granular analysis (strengths vs. weaknesses).",
        "High Inertia: Traditional simulators are expensive, hard to configure, and fail to adapt to specialized job roles (e.g., UX, Data Science) or experience levels."
    ])

    # Slide 4: Proposed Solution
    add_bullets_slide("Proposed Solution: InterviewAce AI", [
        "Interactive Interview Simulation: Tailored specifically to the chosen Job Role and Experience Level.",
        "State-of-the-Art LLM Integration: Uses IBM watsonx.ai running the Granite model for context-aware questions and evaluation.",
        "Granular Scoring Engine: Assigns an objective score (0-10) based on answer completeness, terminology, and structure.",
        "Actionable Feedback: Highlights explicit strengths and weaknesses, alongside a comprehensive suggested answer to aid learning."
    ])

    # Slide 5: System Architecture
    add_bullets_slide("System Architecture", [
        "1. Streamlit Frontend UI: Interactive dashboard for choosing configuration and typing answers.",
        "2. State Manager (Session State): Coordinates history tracker, active questions, and results.",
        "3. Prompt Builder: Interlaces parameters (Role, Experience, Answer) into specialized prompt templates.",
        "4. IBM watsonx.ai Engine: Connects securely via the official Python SDK to process prompts.",
        "5. Granite Model (ibm/granite-3-0-8b-instruct): Evaluates responses and generates next questions.",
        "6. Output Parser: Translates LLM results into JSON formats displayed dynamically as cards on the UI."
    ])

    # Slide 6: Workflow
    add_bullets_slide("Application Workflow", [
        "1. Configuration: Candidate selects Job Role (e.g., Software Engineer) and Experience Level (e.g., Fresher).",
        "2. Generate Question: App fetches a customized technical or behavioral question from IBM watsonx.ai.",
        "3. Answer Submission: Candidate reads the question, drafts their response, and clicks 'Evaluate Answer'.",
        "4. AI Evaluation: The answer is sent alongside context to IBM Granite to extract a score and specific bullet points.",
        "5. Dashboard Display: UI updates instantly with visual score cards, strengths, weaknesses, and a suggested model answer.",
        "6. Next Question: Candidate clicks 'Next Question' to proceed; session state stores performance history."
    ])

    # Slide 7: Technologies Used
    add_bullets_slide("Technology Stack", [
        "• Frontend Framework: Streamlit (Python-based interactive web framework)",
        "• AI Platform: IBM watsonx.ai (Model Inference Client)",
        "• Foundation LLM: IBM Granite Model (ibm/granite-3-0-8b-instruct)",
        "• Configuration: python-dotenv (Secure environment variable management)",
        "• Deployment: Streamlit Community Cloud (Instant hosting & live demo linkage)",
        "• Version Control: GitHub (Collaborative source repository)"
    ])

    # Slide 8: Screenshots & UI Design
    add_bullets_slide("UI/UX Design Highlights", [
        "• Premium Dark Theme: Custom CSS styling with slate, indigo, and emerald green accent panels.",
        "• Response Dashboard: Clean layouts splitting AI Score, Strengths (success indicator), and Weaknesses (alert indicator).",
        "• API Indicator Badge: Visual notification showing 'IBM watsonx.ai (Active)' or 'Sandbox Mode (Fallback)'.",
        "• Progressive Progression: Easy navigation via responsive CTA buttons ('Evaluate Answer', 'Next Question')."
    ])

    # Slide 9: Future Scope
    add_bullets_slide("Future Enhancements", [
        "• Resume Upload Integration: Parse candidate resumes to generate personalized, hyper-focused interview questions.",
        "• Voice Interviewing: Integrate WebRTC / speech-to-text to capture spoken answers, analyzing tone and pacing.",
        "• Company-Specific Mocking: Allow users to specify target companies (e.g., Google, IBM, Amazon) for tailored questions.",
        "• ATS Resume Analyzer: Direct feedback on resumes matching the job description to optimize profile conversion."
    ])

    # Slide 10: Conclusion
    add_bullets_slide("Conclusion & Social Impact", [
        "• Accessibility: Democratizes interview preparation for students and freshers worldwide at zero cost.",
        "• Scalability: Can be easily expanded to support dozens of roles, voice processing, and real-time coding pads.",
        "• Professional Pitch: Satisfies the mandatory IBM Cloud usage requirements with a deployment-ready codebase, live linkage, and solid architecture.",
        "• Repository Linkage: Complete submission package includes clean open-source code on GitHub with comprehensive documentation."
    ])

    # Save presentation
    output_path = "InterviewAce_AI_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
